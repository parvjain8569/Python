from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_number_systems_presentation():
    # Create a presentation object
    prs = Presentation()

    # ========== TITLE SLIDE ==========
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Number Systems and Conversions"
    subtitle.text = "Foundations of Computer Science (25CSE0102)\nDr. Punit Soni\nChitkara University, Punjab"

    # ========== SLIDE 1: Introduction to Number Systems ==========
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Introduction to Number Systems"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Definition: Technique to represent numbers in computer architecture"
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Every value stored or retrieved follows a defined number system"
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Supported systems:"
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• Binary"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Decimal"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Octal"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Hexadecimal"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Foundation for data encoding, arithmetic, and computation"
    p.font.size = Pt(20)

    # ========== SLIDE 2: Binary Number System ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Binary Number System"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    bullets = [
        "Most fundamental system in computers",
        "Uses only two digits: 0 and 1",
        "Based on electronic circuit states: ON (1) and OFF (0)",
        "Base/Radix = 2",
        "Each binary digit = bit",
        "LSB (Least Significant Bit) and MSB (Most Significant Bit)",
        "Example: 11010₂ = 26₁₀"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.space_after = Pt(6)

    # ========== SLIDE 3: Decimal Number System ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Decimal Number System"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    bullets = [
        "Base-10 system",
        "Digits: 0–9",
        "Positional value system",
        "Standard for everyday counting"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)

    # ========== SLIDE 4: Octal Number System ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Octal Number System"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    bullets = [
        "Base-8 system",
        "Digits: 0–7",
        "Positional value system",
        "Used as shorthand for binary (each octal digit = 3 bits)",
        "Example: 467₈ = 311₁₀"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(22)

    # ========== SLIDE 5: Hexadecimal Number System ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Hexadecimal Number System"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    bullets = [
        "Base-16 system",
        "Digits: 0–9, A–F",
        "Widely used in programming",
        "Each hex digit = 4 bits",
        "Example: 27FB₁₆ = 10235₁₀"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(22)

    # ========== SLIDE 6: Number System Conversion Overview ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Number System Conversion Overview"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Purpose: Translate between human-readable and machine-readable formats"
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Common conversions:"
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "• Decimal ↔ Binary"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Binary ↔ Octal"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Binary ↔ Hexadecimal"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Methods: Repeated division/multiplication, positional notation"
    p.font.size = Pt(20)

    # ========== SLIDE 7: Decimal to Binary Conversion ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Decimal to Binary Conversion"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Integral Part:"
    p.font.size = Pt(20)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "1. Divide by 2"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "2. Record remainder (0 or 1)"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "3. Repeat with quotient until quotient = 0"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "4. Write remainders in reverse order"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Example: 278₁₀ = 100010110₂"
    p.font.size = Pt(18)
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "Fractional Part:"
    p.font.size = Pt(20)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "1. Multiply by 2"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "2. Record integer part"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "3. Repeat with fractional part until 0 or desired precision"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Example: 10.25₁₀ = 1010.01₂"
    p.font.size = Pt(18)
    p.font.italic = True

    # ========== SLIDE 8: Binary to Decimal Conversion ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Binary to Decimal Conversion"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Integral Part:"
    p.font.size = Pt(20)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Multiply each bit by 2^position (right to left, starting at 0)"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Sum results"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Example: 101101₂ = 45₁₀"
    p.font.size = Pt(18)
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "Fractional Part:"
    p.font.size = Pt(20)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Multiply each bit by 2^-position (left to right after decimal)"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "• Sum results"
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Example: 0.1011₂ = 0.6875₁₀"
    p.font.size = Pt(18)
    p.font.italic = True

    # ========== SLIDE 9: Decimal to Octal Conversion ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Decimal to Octal Conversion"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    steps = [
        "1. Divide decimal by 8",
        "2. Record remainder",
        "3. Repeat with quotient until quotient = 0",
        "4. Write remainders in reverse order",
        "Example: 92₁₀ = 134₈"
    ]
    
    for step in steps:
        p = tf.add_paragraph()
        if step.startswith("Example"):
            p.text = step
            p.font.italic = True
            p.font.size = Pt(20)
        else:
            p.text = step
            p.font.size = Pt(22)

    # ========== SLIDE 10: Octal to Decimal Conversion ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Octal to Decimal Conversion"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    bullets = [
        "Multiply each digit by 8^position",
        "Sum results",
        "Example: 125₈ = 85₁₀"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)
        if "Example" in bullet:
            p.font.italic = True

    # ========== SLIDE 11: Binary to Octal Conversion ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Binary to Octal Conversion"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    bullets = [
        "Group binary digits into sets of 3 (starting from right)",
        "Replace each group with equivalent octal digit",
        "Example: 10110010101₂ = 2625₈"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)
        if "Example" in bullet:
            p.font.italic = True

    # ========== SLIDE 12: Octal to Binary Conversion ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Octal to Binary Conversion"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    bullets = [
        "Direct method: Convert each octal digit to 3-bit binary",
        "Example: 345₈ = 011100101₂"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(26)
        if "Example" in bullet:
            p.font.italic = True

    # ========== SLIDE 13: Binary to Hexadecimal Conversion ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Binary to Hexadecimal Conversion"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    bullets = [
        "Group binary digits into sets of 4 (quads)",
        "Replace each quad with equivalent hex digit",
        "Example: Provided in practice"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(26)

    # ========== SLIDE 14: Hexadecimal to Binary Conversion ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Hexadecimal to Binary Conversion"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    bullets = [
        "Break hex number into individual digits",
        "Convert each hex digit to 4-bit binary",
        "Combine all binary groups",
        "Example: Provided in practice"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)

    # ========== SLIDE 15: Summary ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Summary"
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    bullets = [
        "Number systems are foundational to computing",
        "Binary is core; octal and hex are compact representations",
        "Conversions enable interoperability between human and machine data"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(28)
        p.space_after = Pt(12)

    # ========== SLIDE 16: Thank You ==========
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You"
    subtitle.text = "Questions?"

    # Save the presentation
    prs.save('Number_Systems_Presentation.pptx')
    print("Presentation created successfully: 'Number_Systems_Presentation.pptx'")

if __name__ == "__main__":
    create_number_systems_presentation()
    