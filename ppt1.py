from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Fundamentals of Computer Science"
subtitle.text = "Prepared by: Dr. Monika Sethi\nChitkara University Institute of Engineering and Technology\nChitkara University, Punjab"

# Slide 2: Introduction to Computers
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Introduction to Computers"
content.text = "• What is a computer?\n• How does it work?\n• Uses in daily life\n• History and generations of computers\n• Computer structure (Von Neumann Architecture)"

# Slide 3: What is a Computer?
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "What is a Computer?"
content.text = "• The word 'computer' comes from 'compute' (to calculate)\n• An electronic device that performs calculations quickly\n• Also called a data processor\n• Can store, process, and get back data when needed"

# Slide 4: Data Processing
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Data Processing"
content.text = "• Data: Raw facts (input)\n• Information: Processed data (output)\n• Steps:\n  1. Capture data\n  2. Work with data\n  3. Show results"

# Slide 5: Computer Characteristics
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Computer Characteristics"
content.text = "• Automatic: Works by itself once started\n• Fast: Works in microseconds/nanoseconds\n• Accurate: Very few errors if input is correct\n• Consistent: Doesn't get tired or bored\n• Versatile: Can do many different tasks\n• Good memory: Stores lots of information\n• No feelings: Follows instructions without emotions"

# Slide 6: Computer Uses - Part 1
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "How We Use Computers"
content.text = "1. Communication: Email, video calls, social media\n2. Learning: Online classes, research tools\n3. Business: Office work, data management\n4. Banking: Online banking, digital payments\n5. Healthcare: Medical records, online check-ups"

# Slide 7: Computer Uses - Part 2
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "More Computer Uses"
content.text = "6. Entertainment: Movies, music, games\n7. Shopping: Online stores, price comparison\n8. Travel: Maps, ride services\n9. Smart Homes: Voice assistants, security\n10. Creative Work: Design, writing, photography"

# Slide 8: History of Computers
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "History of Computers"
content.text = "• 1642: First mechanical calculator (Pascal)\n• 1671: First multiplication calculator (Leibniz)\n• 1880s: Keyboard machines and punched cards\n• 1822-1842: Charles Babbage - Father of modern computers\n• Designed early computer models"

# Slide 9: Computer Generations
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Computer Generations"
content.text = "Generations show technology progress over time:\n\n1. First (1940s-50s): Vacuum tubes, very large\n2. Second (1950s-60s): Transistors, smaller\n3. Third (1960s-70s): Integrated circuits, cheaper\n4. Fourth (1970s-Now): Microprocessors, personal computers\n5. Fifth (Now-Future): AI, smart computers"

# Slide 10: Von Neumann Architecture
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Computer Structure: Von Neumann"
content.text = "• Designed by John von Neumann (1945)\n• Used in most modern computers\n• Key idea: Store programs and data together in memory"

# Slide 11: Computer Parts
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Main Computer Parts"
content.text = "1. Memory: Stores data and instructions\n2. ALU: Does calculations and comparisons\n3. Control Unit: Manages all operations\n4. Input: Keyboard, mouse, scanner\n5. Output: Monitor, printer, speakers\n6. Registers: Fast temporary storage inside CPU"

# Slide 12: How Computers Work
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "How Computers Work"
content.text = "1. Fetch: Get next instruction\n2. Decode: Understand what to do\n3. Execute: Perform the action\n4. Repeat: Go to next instruction\n\nThis happens billions of times per second!"

# Slide 13: Limitations
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Computer Limitations"
content.text = "• Does only what it's told\n• Can't think for itself\n• Needs correct instructions\n• Follows rules exactly\n• No creativity or emotions"

# Save the presentation
prs.save('Computer_Science_Basics.pptx')

print("Presentation created successfully: Computer_Science_Basics.pptx")